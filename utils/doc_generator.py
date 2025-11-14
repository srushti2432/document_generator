import json
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from utils.gemini_api import generate_content

def create_docx(content, path):
    doc = Document()
    doc.add_paragraph(content)
    doc.save(path)

def create_ppt(content, path):
    prs = Presentation()
    for section in content.strip().split("\n\n"):
        lines = section.split("\n")
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = lines[0] if lines else "Untitled"
        slide.placeholders[1].text = "\n".join(lines[1:])
    prs.save(path)

def create_excel(content, path):
    wb = Workbook()
    ws = wb.active
    for row in content.strip().split("\n"):
        ws.append([cell.strip() for cell in row.split(",")])
    wb.save(path)

def generate_brd(requirements, prompt_path):
    with open(prompt_path, "r") as f:
        prompt = f.read().replace("{{requirements}}", requirements)
    content = generate_content(prompt)
    return content.encode("utf-8")

def generate_hld(requirements, prompt_path):
    with open(prompt_path, "r") as f:
        prompt = f.read().replace("{{requirements}}", requirements)
    content = generate_content(prompt)
    return content.encode("utf-8")

def generate_fsd(parsed_code, requirements, prompt_path):
    formatted_code_summary = json.dumps(parsed_code, indent=2)
    with open(prompt_path, "r") as f:
        prompt = f.read().replace("{{code_summary}}", formatted_code_summary).replace("{{requirements}}", requirements)
    content = generate_content(prompt)
    return content.encode("utf-8")

def generate_lld(parsed_code, requirements, prompt_path):
    formatted_code_summary = json.dumps(parsed_code, indent=2)
    with open(prompt_path, "r") as f:
        prompt = f.read().replace("{{code_summary}}", formatted_code_summary).replace("{{requirements}}", requirements)
    content = generate_content(prompt)
    return content.encode("utf-8")

def generate_ppt(parsed_code, requirements, prompt_path):
    formatted_code_summary = json.dumps(parsed_code, indent=2)
    with open(prompt_path, "r") as f:
        prompt = f.read().replace("{{code_summary}}", formatted_code_summary).replace("{{requirements}}", requirements)
    content = generate_content(prompt)
    return content.encode("utf-8")

def generate_excel(requirements, parsed_code, prompt_path):
    formatted_code_summary = json.dumps(parsed_code, indent=2)
    with open(prompt_path, "r") as f:
        prompt = f.read().replace("{{requirements}}", requirements).replace("{{code_summary}}", formatted_code_summary)
    content = generate_content(prompt)
    return content.encode("utf-8")
